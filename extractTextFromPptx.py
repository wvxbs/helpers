from pptx import Presentation
from dotenv import load_dotenv
import os

def extractTextFromPptx(inputFile):
    prs = Presentation(inputFile)
    textContent = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                textContent.append(shape.text)
    return "\n".join(textContent)

load_dotenv()

pptxFile = os.getenv("PPTX_FILE")
outputFile = os.getenv("PPTX_OUTPUT_FILE")

plainText = extractTextFromPptx(pptxFile)

with open(outputFile, "w", encoding="utf-8") as f:
    f.write(plainText)

print("Exportação concluída!")